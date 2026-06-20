import os
import fitz
import unicodedata
import re
import time
import json
import html
from PIL import Image
import pdfplumber
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tenacity import retry, wait_exponential, stop_after_attempt, retry_if_exception
from transformers import AutoTokenizer
import underthesea
from qdrant_client import models
from logger_config import logger
# FIX #1: dung 2 ham moi (reset 1 lan + insert tung trang). Giu save_document_metadata cho file 1 trang.
from db_logic import reset_document_metadata, save_page_metadata, save_document_metadata, save_bom_records, get_document_info
from rag_logic import vectorstore, client
# FIX Bug #4: predicate retry dung chung cho Gemini (google-genai)
from gemini_client import describe_gemini_error, is_retryable_error
from functools import lru_cache

def _env_bool(name, default=False):
    raw = os.getenv(name)
    if raw is None:
        return default
    return str(raw).strip().lower() in {"1", "true", "yes", "y", "on"}

EMBEDDING_CHUNK_SIZE = int(os.getenv("EMBEDDING_CHUNK_SIZE", "220"))
EMBEDDING_CHUNK_OVERLAP = int(os.getenv("EMBEDDING_CHUNK_OVERLAP", "40"))
STRICT_INGEST_REQUIRE_VISION = _env_bool("STRICT_INGEST_REQUIRE_VISION", True)
ROLLBACK_ON_INGEST_ERROR = _env_bool("ROLLBACK_ON_INGEST_ERROR", True)
GEMINI_METADATA_MODE = os.getenv("GEMINI_METADATA_MODE", "missing_only").strip().lower()

def remove_accents(text: str) -> str:
    if text is None:
        return ""
    text = str(text)
    text = text.replace("đ", "d").replace("Đ", "D")
    normalized = unicodedata.normalize("NFD", text)
    return "".join(
        ch for ch in normalized
        if unicodedata.category(ch) != "Mn"
    )
 
# FIX H5: Cache word_tokenize (underthesea cham 50-200ms/call). Header / chunk lap lai khong tokenize lai.
@lru_cache(maxsize=4096)
def tokenize_cached(text):
    return underthesea.word_tokenize(text, format="text")
 
# Khoi tao Tokenizer Singleton de do luong chunk size chinh xac
EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL", "BAAI/bge-m3")

try:
    GLOBAL_TOKENIZER = AutoTokenizer.from_pretrained(EMBEDDING_MODEL_NAME)
    logger.info(f"Da load AutoTokenizer ({EMBEDDING_MODEL_NAME}) thanh cong cho Chunking.")
except Exception as e:
    GLOBAL_TOKENIZER = None
    logger.warning(f"Khong load duoc tokenizer {EMBEDDING_MODEL_NAME}: {e}")
 
def tokenizer_length(text):
    if GLOBAL_TOKENIZER:
        return len(GLOBAL_TOKENIZER.encode(text))
    return len(text)
 
# Token-based Text Splitter dung chung cho toan bo module
token_splitter = RecursiveCharacterTextSplitter(
    chunk_size=EMBEDDING_CHUNK_SIZE,
    chunk_overlap=EMBEDDING_CHUNK_OVERLAP,
    length_function=tokenizer_length
)
 
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMAGE_DIR = os.path.join(BASE_DIR, "Data_Anh_Da_Tach")
os.makedirs(IMAGE_DIR, exist_ok=True)
 
try:
    import pandas as pd
except ImportError:
    pd = None
 
try:
    import docx
except ImportError:
    docx = None
 
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None
 
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
 
PDF_EXTENSIONS = {".pdf"}
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".sql",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".cs",
    ".cpp", ".c", ".h", ".hpp", ".json", ".xml", ".yaml",
    ".yml", ".ini", ".cfg",
}
HTML_EXTENSIONS = {".html", ".htm"}
TABLE_EXTENSIONS = {".csv", ".tsv", ".xlsx", ".xls"}
WORD_EXTENSIONS = {".docx"}
PRESENTATION_EXTENSIONS = {".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff"}
 
SUPPORTED_LEARNING_EXTENSIONS = (
    PDF_EXTENSIONS
    | TEXT_EXTENSIONS
    | HTML_EXTENSIONS
    | TABLE_EXTENSIONS
    | WORD_EXTENSIONS
    | PRESENTATION_EXTENSIONS
    | IMAGE_EXTENSIONS
)
 
# 1. Ham trich xuat thong minh (Di chuyen tu file cu sang)
def _metadata_needs_llm(result):
    if GEMINI_METADATA_MODE in {"off", "false", "0", "none"}:
        return False
    if GEMINI_METADATA_MODE == "always":
        return True
    
    if not result.get("ma_doi_tuong"):
        return True
        
    # Neu co nhieu ma nhung chua phan loai duoc BTP/vat tu thi nen goi LLM
    if len(result.get("ma_doi_tuong", [])) >= 2 and not result.get("ma_btp") and not result.get("ma_vat_tu"):
        return True
        
    critical_fields = ("ten_tai_lieu", "loai_tai_lieu", "vat_lieu")
    return any(str(result.get(field) or "").strip() in {"", "Khong ro"} for field in critical_fields)

def extract_metadata_smart(text, ten_file, thu_muc, vision_model=None, quality_warnings=None):
    """
    Chien luoc: Regex-first -> Gemini-fallback.
    Luu y: LLM chi duoc goi khi Regex tra ve "Khong ro".
    Regex sai (false positive) se khong duoc LLM tu dong sua.
    Neu format ban ve thay doi dot ngot, kiem tra Regex truoc.
    """
    lines = text.split('\n')
 
    ma_doi_tuong_regex = []
    code_patterns = [
        r"\b\d{1,2}\.\d{1,2}\.\d{3,6}(?:\.\d{1,4})?\b",
        r"\b[A-Z]{1,5}[-_/]?\d{3,8}(?:[-_/][A-Z0-9]+)?\b",
    ]
    
    for pat in code_patterns:
        for m in re.findall(pat, ten_file):
            if m not in ma_doi_tuong_regex:
                ma_doi_tuong_regex.append(m)

    ten_sp_val = "Khong ro"
    vat_lieu_val = "Khong ro"
    
    for idx, line in enumerate(lines):
        line_stripped = line.strip()
        
        for pat in code_patterns:
            for m in re.findall(pat, line_stripped):
                if m not in ma_doi_tuong_regex:
                    ma_doi_tuong_regex.append(m)
                    if ten_sp_val == "Khong ro" and idx + 1 < len(lines):
                        ten_sp_parts = []
                        for j in range(idx + 1, min(idx + 3, len(lines))):
                            next_line = lines[j].strip()
                            if next_line.startswith("Ban ve") or next_line == "" or next_line == "-":
                                break
                            ten_sp_parts.append(next_line)
                        if ten_sp_parts:
                            ten_sp_val = " ".join(ten_sp_parts)
                            
        if re.match(r'^(?:Inox|SUS|SS|AL|Thep|SPCC|Q235|Nhom|Dong|Sat)', line_stripped, re.IGNORECASE):
            if vat_lieu_val == "Khong ro":
                vat_lieu_val = line_stripped

    so_luong_val = "Khong ro"
    found_vat_lieu = False
    for idx, line in enumerate(lines):
        line_stripped = line.strip()
        if re.match(r'^(?:Inox|SUS|SS|AL|Thep|SPCC|Q235|Nhom|Dong|Sat)', line_stripped, re.IGNORECASE):
            found_vat_lieu = True
            continue
        if found_vat_lieu and re.match(r'^\d{1,3}$', line_stripped):
            so_luong_val = line_stripped
            break
        if found_vat_lieu and line_stripped:
            break
 
    cong_doan_val = "Khong ro"
    cong_doan_match = re.search(r'Ban ve\s+(To\s+[\w\s]+?)(?:\n|$)', text)
    if cong_doan_match:
        cong_doan_val = cong_doan_match.group(1).strip()
    else:
        folder_map = {
            "To_Han": "To han", "To_Nham": "To nham", "To_Son": "To son",
            "To_Dong_Goi": "To dong goi", "To_Tien_Phay": "To Tien Phay",
        }
        cong_doan_val = folder_map.get(thu_muc, thu_muc)
 
    ngay_ve_val = "Khong ro"
    ngay_ve_match = re.search(r'(\d{2}/\d{2}/\d{4})', text)
    if ngay_ve_match:
        ngay_ve_val = ngay_ve_match.group(1)
 
    nguoi_lap_val = "Khong ro"
    for idx, line in enumerate(lines):
        if "Ten san pham" in line and ":" in line:
            if idx + 1 < len(lines):
                next_line = lines[idx + 1].strip()
                if next_line and next_line != "Ma BTP :" and not next_line.startswith("CONG TY"):
                    nguoi_lap_val = next_line
            break
 
    dung_sai_day = "Khong ro"
    dung_sai_khac = "Khong ro"
    ds_day_match = re.search(r'Dung sai do day vat lieu\s*:\s*([^\n]+)', text)
    if ds_day_match:
        dung_sai_day = ds_day_match.group(1).strip()
    ds_khac_match = re.search(r'Dung sai cac kich thuoc khac\s*:\s*([^\n]+)', text)
    if ds_khac_match:
        dung_sai_khac = ds_khac_match.group(1).strip()
 
    yckt_text = ""
    in_yckt = False
    for line in lines:
        line_stripped = line.strip()
        if line_stripped.startswith("Ban ve To"):
            in_yckt = True
            continue
        if in_yckt:
            if re.match(r'^9\.3\.\d{4,5}$', line_stripped):
                break
            if line_stripped.startswith("-") and len(line_stripped) > 1:
                yckt_text += line_stripped[1:].strip() + "\n"
            elif line_stripped and line_stripped != "-":
                if yckt_text and not yckt_text.endswith("\n"):
                    yckt_text += " " + line_stripped
                elif yckt_text:
                    yckt_text = yckt_text.rstrip("\n") + " " + line_stripped + "\n"
 
    hdcv_val = ""
    all_hdcv = re.findall(r'HDCV:\s*([^\n]+)', text)
    if len(all_hdcv) > 0:
        hdcv_val = " | ".join(all_hdcv)
 
    kich_thuoc_val = ""
    kt_match = re.search(r'Kich thuoc tong the\s*:\s*([^\n]+)', text)
    if kt_match:
        kich_thuoc_val = kt_match.group(1).strip()
    elif re.search(r'(\d{2,4}\s*[xX]\s*\d{2,4}\s*[xX]\s*\d{2,4})\s*mm', text):
        kich_thuoc_val = re.search(r'(\d{2,4}\s*[xX]\s*\d{2,4}\s*[xX]\s*\d{2,4})\s*mm', text).group(1) + "mm"

    ma_chinh_regex = []
    ma_lien_quan_regex = []
    
    for m in ma_doi_tuong_regex:
        if m in ten_file:
            ma_chinh_regex.append(m)
        else:
            ma_lien_quan_regex.append(m)
            
    if not ma_chinh_regex and ma_doi_tuong_regex:
        ma_chinh_regex = [ma_doi_tuong_regex[0]]
        ma_lien_quan_regex = ma_doi_tuong_regex[1:]

    result = {
        "ma_doi_tuong": ma_doi_tuong_regex,
        "ma_chinh": ma_chinh_regex,
        "ma_btp": [],
        "ma_vat_tu": [],
        "ma_lien_quan": ma_lien_quan_regex,
        "ten_tai_lieu": ten_sp_val,
        "loai_tai_lieu": "Ban ve gia cong",  # Default
        "cong_doan": cong_doan_val, "vat_lieu": vat_lieu_val, "so_luong": so_luong_val,
        "nguoi_lap": nguoi_lap_val, "ngay_ve": ngay_ve_val, "dung_sai_day": dung_sai_day,
        "dung_sai_khac": dung_sai_khac, "kich_thuoc": kich_thuoc_val,
        "yckt": yckt_text.strip(), "hdcv": hdcv_val
    }
 
    # HYBRID APPROACH: LLM Extraction de doc moi ma (V2).
    # Mac dinh chi goi Gemini khi metadata quan trong con thieu de giam rate limit.
    if vision_model and _metadata_needs_llm(result):
        prompt = f"""
        Ban la chuyen gia doc tai lieu co khi. Hay trich xuat cac thong tin sau tu doan text, tra ve dung dinh dang JSON:
            "ma_chinh": ["ma 1"],
            "ma_btp": ["ma 2"],
            "ma_vat_tu": ["ma 3"],
            "ma_lien_quan": ["ma 4"],
            "ten_tai_lieu": "Ten san pham hoac tieu de tai lieu",
            "loai_tai_lieu": "Nhan ngan gon mo ta tai lieu (VD: Ban ve gia cong, So tay ISO, Catalog...)",
            "vat_lieu": "Vat lieu de cap (neu co)"
            
        Goi y cac thong tin so bo da tim thay (hay kiem tra, mo rong hoac sua lai neu can):
        - Ma chinh: {result.get("ma_chinh", [])}
        - Ma lien quan: {result.get("ma_lien_quan", [])}
        - Ten tai lieu: {result.get("ten_tai_lieu")}
        - Vat lieu: {result.get("vat_lieu")}
        
        Uu tien ket qua phan tich cua ban neu hop ly hon. Luu y: Chi tra ve dung JSON, khong giai thich gi them.
        Text can phan tich:
        {text}
        """
        try:
            import json
            response = call_gemini_vision(vision_model, prompt)
            raw_json = response.text.replace('```json', '').replace('```', '').strip()
            llm_result = json.loads(raw_json)
            
            for key in ["ma_chinh", "ma_btp", "ma_vat_tu", "ma_lien_quan"]:
                if key in llm_result and isinstance(llm_result[key], list) and llm_result[key]:
                    result[key] = [str(x) for x in llm_result[key]]
                    
            # Combine all codes to ma_doi_tuong for backward compatibility
            all_codes = result["ma_chinh"] + result["ma_btp"] + result["ma_vat_tu"] + result["ma_lien_quan"]
            result["ma_doi_tuong"] = list(dict.fromkeys(all_codes))
            
            if "ten_tai_lieu" in llm_result and llm_result["ten_tai_lieu"]:
                result["ten_tai_lieu"] = str(llm_result["ten_tai_lieu"])
            if "loai_tai_lieu" in llm_result and llm_result["loai_tai_lieu"]:
                result["loai_tai_lieu"] = str(llm_result["loai_tai_lieu"])
            if "vat_lieu" in llm_result and llm_result["vat_lieu"]:
                result["vat_lieu"] = str(llm_result["vat_lieu"])
            if "quality_warnings" in llm_result and isinstance(llm_result["quality_warnings"], list) and quality_warnings is not None:
                for w in llm_result["quality_warnings"]:
                    if w not in quality_warnings:
                        quality_warnings.append(str(w))
        except Exception as e:
            detail = describe_gemini_error(e)
            msg = f"Loi LLM Fallback boc tach metadata cho {ten_file}: {detail}"
            logger.error(msg)
            if quality_warnings is not None:
                quality_warnings.append(msg)
 
    return result
 
# 2. Co che Retry thong minh cho Gemini chong 429 (da migrate sang google-genai)
@retry(
    retry=retry_if_exception(is_retryable_error),
    wait=wait_exponential(multiplier=2, min=4, max=60),
    stop=stop_after_attempt(5)
)
def call_gemini_vision(vision_model, prompt, image=None):
    if image:
        return vision_model.generate_content([prompt, image])
    else:
        return vision_model.generate_content(prompt)

@retry(
    wait=wait_exponential(multiplier=2, min=2, max=30),
    stop=stop_after_attempt(4),
    reraise=True
)
def _add_docs_with_retry(chunks):
    vectorstore.add_documents(chunks)

def _delete_vectors_for_file(ten_file, thu_muc):
    client.delete(
        collection_name="TaiLieuKyThuat_v2",
        points_selector=models.Filter(must=[
            models.FieldCondition(key="metadata.file_goc", match=models.MatchValue(value=ten_file)),
            models.FieldCondition(key="metadata.phong_ban_quyen", match=models.MatchValue(value=thu_muc)),
        ])
    )
 
def _require_package(package, feature_name):
    if package is None:
        raise ImportError(f"Thieu thu vien de doc {feature_name}. Hay cai/cap nhat requirements.txt roi chay lai.")
 
def _read_text_file(file_path):
    encodings = ("utf-8-sig", "utf-8", "cp1258", "cp1252", "latin-1")
    last_error = None
    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError as exc:
            last_error = exc
    raise UnicodeDecodeError(
        last_error.encoding,
        last_error.object,
        last_error.start,
        last_error.end,
        "Khong doc duoc file bang cac encoding pho bien.",
    )
 
def _read_json_file(file_path):
    raw_text = _read_text_file(file_path)
    try:
        data = json.loads(raw_text)
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception:
        return raw_text
 
def _read_xml_file(file_path):
    raw_text = _read_text_file(file_path)
    try:
        import xml.etree.ElementTree as ET
        root = ET.fromstring(raw_text)
        extracted = "\n".join(t.strip() for t in root.itertext() if t and t.strip())
        return extracted or raw_text
    except Exception:
        return raw_text
 
def _read_html_file(file_path):
    raw_text = _read_text_file(file_path)
    if BeautifulSoup:
        soup = BeautifulSoup(raw_text, "html.parser")
        return soup.get_text("\n")
    without_tags = re.sub(r"<[^>]+>", " ", raw_text)
    return html.unescape(without_tags)
 
def _dataframe_to_text(df, title=None):
    if df is None or df.empty:
        return f"{title}\n(Bang rong)" if title else "(Bang rong)"
    df = df.fillna("")
    text = df.to_string(index=False)
    return f"{title}\n{text}" if title else text
 
def _read_table_file(file_path, ext):
    _require_package(pd, "bang tinh/CSV")
    if ext in {".csv", ".tsv"}:
        sep = "\t" if ext == ".tsv" else ","
        last_error = None
        for encoding in ("utf-8-sig", "utf-8", "cp1258", "cp1252", "latin-1"):
            try:
                df = pd.read_csv(file_path, sep=sep, encoding=encoding)
                return _dataframe_to_text(df, title=f"Bang du lieu tu {os.path.basename(file_path)}")
            except UnicodeDecodeError as exc:
                last_error = exc
        raise last_error
    sheets = pd.read_excel(file_path, sheet_name=None)
    sheet_texts = []
    for sheet_name, df in sheets.items():
        sheet_texts.append(_dataframe_to_text(df, title=f"Sheet: {sheet_name}"))
    return "\n\n---\n\n".join(sheet_texts)
 
def _read_word_file(file_path):
    _require_package(docx, "Word DOCX")
    document = docx.Document(file_path)
    parts = [para.text.strip() for para in document.paragraphs if para.text and para.text.strip()]
    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"Bang {table_index}:\n" + "\n".join(rows))
    return "\n\n".join(parts)
 
def _read_presentation_file(file_path):
    _require_package(Presentation, "PowerPoint PPTX")
    prs = Presentation(file_path)
    slides = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    parts.append("Bang:\n" + "\n".join(rows))
        if parts:
            slides.append(f"Slide {slide_index}:\n" + "\n\n".join(parts))
    return "\n\n---\n\n".join(slides)
 
def parse_vision_json(raw_text):
    """Parse JSON tu Gemini tra ve, ho tro fallback regex"""
    try:
        json_str = raw_text
        if "```json" in json_str:
            json_str = json_str.split("```json")[1].split("```")[0].strip()
        elif "```" in json_str:
            json_str = json_str.split("```")[1].split("```")[0].strip()
        else:
            match = re.search(r'\{.*\}', json_str, re.DOTALL)
            if match:
                json_str = match.group(0)
                
        data = json.loads(json_str)
        return data
    except Exception as e:
        logger.warning(f"Khong the parse JSON tu vision: {e}")
        return None

def format_vision_data(data):
    """Format dictionary thanh text de dua vao Qdrant"""
    if not data:
        return ""
    
    parts = []
    if data.get("document_codes"):
        parts.append(f"- Mã bản vẽ/tài liệu: {', '.join([str(x) for x in data['document_codes']])}")
    if data.get("part_names"):
        parts.append(f"- Tên chi tiết/vật tư: {', '.join([str(x) for x in data['part_names']])}")
    if data.get("materials"):
        parts.append(f"- Vật liệu: {', '.join([str(x) for x in data['materials']])}")
    if data.get("dimensions"):
        parts.append(f"- Kích thước nổi bật: {', '.join([str(x) for x in data['dimensions']])}")
    if data.get("tolerances"):
        parts.append(f"- Dung sai: {', '.join([str(x) for x in data['tolerances']])}")
    if data.get("technical_notes"):
        parts.append("- Ghi chú kỹ thuật:\n  " + "\n  ".join([str(x) for x in data["technical_notes"]]))
    if data.get("bom_rows"):
        parts.append("- Các hàng vật tư (BOM) nhận diện được:\n  " + "\n  ".join([str(row) for row in data["bom_rows"]]))
    if data.get("uncertain_fields"):
        parts.append(f"- Các phần mờ/không chắc chắn: {', '.join([str(x) for x in data['uncertain_fields']])}")
        
    return "Kết quả phân tích ảnh (Cấu trúc JSON):\n" + "\n".join(parts) if parts else ""
def _read_image_file(file_path, ten_file, vision_model):
    if not vision_model:
        raise ValueError("File anh can GOOGLE_API_KEY hop le de Gemini Vision doc noi dung/OCR.")
    image = Image.open(file_path)
    prompt = (
        f"Day la file anh '{ten_file}' duoc nap lam du lieu cho chatbot ky thuat. "
        "Hay OCR va tra ve ket qua DUOI DANG JSON voi schema sau:\n"
        "{\n"
        '  "document_codes": [],\n'
        '  "part_names": [],\n'
        '  "materials": [],\n'
        '  "dimensions": [],\n'
        '  "tolerances": [],\n'
        '  "technical_notes": [],\n'
        '  "bom_rows": [],\n'
        '  "uncertain_fields": []\n'
        "}\n"
        "Luon tra ve dung dinh dang JSON (khong kem text mo dau/ket thuc ngoai block ```json). "
        "Dien vao cac mang cac thong tin ky thuat tuong ung ban nhin thay trong anh."
    )
    response = call_gemini_vision(vision_model, prompt, image)
    return format_vision_json(response.text)
 
def extract_text_from_supported_file(file_path, ten_file, vision_model=None):
    ext = os.path.splitext(file_path)[1].lower()
    if ext in {".json"}:
        return _read_json_file(file_path), "du_lieu_json"
    if ext in {".xml"}:
        return _read_xml_file(file_path), "du_lieu_xml"
    if ext in HTML_EXTENSIONS:
        return _read_html_file(file_path), "van_ban_html"
    if ext in TABLE_EXTENSIONS:
        return _read_table_file(file_path, ext), "bang_du_lieu"
    if ext in WORD_EXTENSIONS:
        return _read_word_file(file_path), "van_ban_word"
    if ext in PRESENTATION_EXTENSIONS:
        return _read_presentation_file(file_path), "slide"
    if ext in IMAGE_EXTENSIONS:
        return _read_image_file(file_path, ten_file, vision_model), "image_summary"
    if ext in TEXT_EXTENSIONS:
        return _read_text_file(file_path), "van_ban"
    supported = ", ".join(sorted(SUPPORTED_LEARNING_EXTENSIONS))
    raise ValueError(f"Dinh dang {ext or '(khong co duoi file)'} chua duoc ho tro. Cac dinh dang dang ho tro: {supported}")
 
def extract_bom_records(table, table_idx=None):
    records = []
    if not table or len(table) < 2: return records
    
    cleaned_table = []
    for row in table:
        cleaned_row = [str(cell).replace("\n", " ").strip() if cell else "" for cell in row]
        cleaned_table.append(cleaned_row)
        
    header_idx = -1
    col_map = {'ma': -1, 'ten': -1, 'vat_lieu': -1, 'sl': -1, 'ghi_chu': -1, 'unit': -1}
    
    for row_idx in range(min(5, len(cleaned_table))):
        row_norm = [remove_accents(h.lower()) for h in cleaned_table[row_idx]]
        
        has_ma_or_ten = any(kw in h for h in row_norm for kw in ['ma hang', 'ma vat tu', 'ma chi tiet', 'ma btp', 'ma tp', 'ky hieu', 'ten vat tu', 'vat tu', 'mo ta', 'ten hang', 'chi tiet', 'ten goi'])
        has_sl_or_vatlieu = any(kw in h for h in row_norm for kw in ['so luong', 'sl', 'vat lieu'])
        
        if has_ma_or_ten and has_sl_or_vatlieu:
            header_idx = row_idx
            for i, h in enumerate(row_norm):
                if any(kw in h for kw in ['ma hang', 'ma vat tu', 'ma chi tiet', 'ma btp', 'ma tp', 'ky hieu']) and col_map['ma'] == -1: col_map['ma'] = i
                elif any(kw in h for kw in ['ten vat tu', 'vat tu', 'mo ta', 'ten hang', 'chi tiet', 'ten goi']) and col_map['ten'] == -1: col_map['ten'] = i
                elif 'vat lieu' in h and col_map['vat_lieu'] == -1: col_map['vat_lieu'] = i
                elif ('so luong' in h or h == 'sl') and col_map['sl'] == -1: col_map['sl'] = i
                elif 'ghi chu' in h and col_map['ghi_chu'] == -1: col_map['ghi_chu'] = i
                elif any(kw in h for kw in ['don vi', 'dvt', 'unit']) and col_map['unit'] == -1: col_map['unit'] = i
            break
            
    if header_idx != -1:
        import json
        for row in cleaned_table[header_idx + 1:]:
            rec = {}
            if col_map['ma'] != -1 and col_map['ma'] < len(row): rec['ma_hang'] = row[col_map['ma']]
            if col_map['ten'] != -1 and col_map['ten'] < len(row): rec['ten_vat_tu'] = row[col_map['ten']]
            if col_map['vat_lieu'] != -1 and col_map['vat_lieu'] < len(row): rec['vat_lieu'] = row[col_map['vat_lieu']]
            if col_map['sl'] != -1 and col_map['sl'] < len(row): 
                try: 
                    num_str = re.sub(r'\D', '', row[col_map['sl']])
                    rec['so_luong'] = int(num_str) if num_str else None
                except: rec['so_luong'] = None
            if col_map['ghi_chu'] != -1 and col_map['ghi_chu'] < len(row): rec['ghi_chu'] = row[col_map['ghi_chu']]
            if col_map['unit'] != -1 and col_map['unit'] < len(row): rec['don_vi'] = row[col_map['unit']]
            
            rec['confidence'] = 0.9 if rec.get('ma_hang') and rec.get('so_luong') else 0.5
            rec['raw_row_json'] = json.dumps(row, ensure_ascii=False)
            rec['source_table_index'] = table_idx
            
            if rec.get('ten_vat_tu') or rec.get('ma_hang'):
                records.append(rec)
    return records

# 3. Ham xu ly PDF trung tam
def process_and_ingest_pdf(pdf_path, ten_file, thu_muc, vision_model=None, progress_callback=None):
    start_time = time.time()
    report = {
        "status": "success",
        "ten_file": ten_file,
        "total_pages": 0,
        "total_chunks": 0,
        "failed_pages": [],
        "vision_failed_pages": [],
        "metadata_llm_failed_pages": [],
        "warnings": [],
        "time_taken": 0,
        "message": ""
    }
    doc = None
    pdf_table_reader = None
    try:
        doc = fitz.open(pdf_path)
        report["total_pages"] = len(doc)
 
        # FIX #1: Reset metadata MOT LAN cho ca file, lay doc_id dung chung cho moi trang
        doc_id = reset_document_metadata(ten_file, thu_muc)
        doc_info = get_document_info(doc_id)
 
        # FIX hieu nang: mo pdfplumber MOT LAN ngoai vong lap (truoc day mo lai moi trang)
        try:
            pdf_table_reader = pdfplumber.open(pdf_path)
        except Exception as e:
            logger.warning(f"Khong mo duoc pdfplumber cho {ten_file}: {e}")
            pdf_table_reader = None
 
        base_name = os.path.splitext(ten_file)[0]  # FIX: thay ten_file.replace('.pdf', '')
 
        for page_num in range(len(doc)):
            if progress_callback:
                progress_callback(f"Dang xu ly trang {page_num+1}/{len(doc)}")
            try:
                page = doc.load_page(page_num)
                text = page.get_text("text")

                # Render image truoc de kip phan tich
                pix = page.get_pixmap(dpi=200)
                safe_thu_muc = re.sub(r'[\\/*?:"<>|]', "", thu_muc) if thu_muc else ""
                if safe_thu_muc:
                    img_name = f"{safe_thu_muc}_{base_name}_page{page_num+1}.png"
                else:
                    img_name = f"{base_name}_page{page_num+1}.png"
                img_path = os.path.join(IMAGE_DIR, img_name)
                pix.save(img_path)

                image_summary = ""
                vision_metadata = {}
                is_text_heavy = len(text.strip()) > 1500  # Chi bo qua trang van ban thuan
                
                # Luon goi Gemini cho PDF ky thuat, tru trang toan text
                vision_required = os.path.exists(img_path) and not is_text_heavy
                vision_failed = False
                if vision_model and vision_required:
                    if progress_callback:
                        progress_callback(f"Dang dung AI (Gemini) phan tich anh trang {page_num+1}...")
                    try:
                        img_to_analyze = Image.open(img_path)
                        prompt = (
                            f"Day la trang so {page_num+1} cua file {ten_file}. "
                            "Hay OCR va tra ve ket qua DUOI DANG JSON voi schema sau:\n"
                            "{\n"
                            '  "document_codes": [],\n'
                            '  "part_names": [],\n'
                            '  "materials": [],\n'
                            '  "dimensions": [],\n'
                            '  "tolerances": [],\n'
                            '  "technical_notes": [],\n'
                            '  "bom_rows": [],\n'
                            '  "uncertain_fields": []\n'
                            "}\n"
                            "Luon tra ve dung dinh dang JSON (khong kem text mo dau/ket thuc ngoai block ```json). "
                            "Dien vao cac mang cac thong tin ky thuat tuong ung ban nhin thay trong hinh."
                        )
                        response = call_gemini_vision(vision_model, prompt, img_to_analyze)
                        vision_data = parse_vision_json(response.text)
                        
                        if vision_data:
                            image_summary = format_vision_data(vision_data)
                            
                            # Add to vision_metadata
                            if vision_data.get("document_codes"): vision_metadata["vision_document_codes"] = ", ".join([str(x) for x in vision_data["document_codes"]])
                            if vision_data.get("part_names"): vision_metadata["vision_part_names"] = ", ".join([str(x) for x in vision_data["part_names"]])
                            if vision_data.get("materials"): vision_metadata["vision_materials"] = ", ".join([str(x) for x in vision_data["materials"]])
                            if vision_data.get("dimensions"): vision_metadata["vision_dimensions"] = ", ".join([str(x) for x in vision_data["dimensions"]])
                            if vision_data.get("tolerances"): vision_metadata["vision_tolerances"] = ", ".join([str(x) for x in vision_data["tolerances"]])
                            if vision_data.get("technical_notes"): vision_metadata["vision_technical_notes"] = ", ".join([str(x) for x in vision_data["technical_notes"]])
                            if vision_data.get("uncertain_fields"): vision_metadata["vision_uncertain_fields"] = ", ".join([str(x) for x in vision_data["uncertain_fields"]])
                            
                            # Convert BOM rows if present
                            if vision_data.get("bom_rows"):
                                structured_bom = []
                                for idx, row in enumerate(vision_data["bom_rows"]):
                                    if isinstance(row, dict):
                                        structured_bom.append({
                                            "ma_hang": str(row.get("ma", row.get("code", row.get("ma_hang", "")))),
                                            "ten_vat_tu": str(row.get("ten", row.get("name", row.get("ten_vat_tu", "")))),
                                            "vat_lieu": str(row.get("vat_lieu", row.get("material", ""))),
                                            "so_luong": row.get("sl", row.get("qty", row.get("so_luong", None))),
                                            "ghi_chu": str(row.get("ghi_chu", row.get("note", ""))),
                                            "don_vi": str(row.get("don_vi", row.get("unit", ""))),
                                            "confidence": row.get("confidence", 0.85),
                                            "raw_row_json": json.dumps(row, ensure_ascii=False),
                                            "source_table_index": 0
                                        })
                                if structured_bom:
                                    save_bom_records(doc_id, page_num + 1, structured_bom)
                        else:
                            image_summary = response.text
                    except Exception as e:
                        vision_failed = True
                        detail = describe_gemini_error(e)
                        warn = f"Trang {page_num+1}: Gemini Vision/OCR loi cho {img_name}: {detail}"
                        report["vision_failed_pages"].append(page_num+1)
                        report["warnings"].append(warn)
                        logger.error(warn)
                elif vision_required and not vision_model:
                    vision_failed = True
                    warn = f"Trang {page_num+1}: can Gemini Vision/OCR nhung chua cau hinh GOOGLE_API_KEY hop le."
                    report["vision_failed_pages"].append(page_num+1)
                    report["warnings"].append(warn)
                    logger.error(warn)

                if vision_failed and STRICT_INGEST_REQUIRE_VISION:
                    report["failed_pages"].append(page_num+1)
                    logger.error(
                        f"Bo qua nap trang {page_num+1} cua {ten_file} de tranh nap thieu du lieu hinh anh/OCR."
                    )
                    continue

                combined_text_for_metadata = text + "\n\n" + image_summary
                warning_count_before_metadata = len(report["warnings"])
                info = extract_metadata_smart(
                    combined_text_for_metadata,
                    ten_file,
                    thu_muc,
                    vision_model,
                    quality_warnings=report["warnings"],
                )
                if len(report["warnings"]) > warning_count_before_metadata:
                    report["metadata_llm_failed_pages"].append(page_num+1)

                metadata = {
                    "file_goc": ten_file,
                    "phong_ban_quyen": thu_muc,
                    "ma_doi_tuong": info["ma_doi_tuong"],
                    "ma_chinh": info.get("ma_chinh", []),
                    "ma_btp": info.get("ma_btp", []),
                    "ma_vat_tu": info.get("ma_vat_tu", []),
                    "ma_lien_quan": info.get("ma_lien_quan", []),
                    "loai_tai_lieu": info["loai_tai_lieu"],
                    "ten_san_pham": info["ten_tai_lieu"],
                    "cong_doan": info["cong_doan"],
                    "so_luong": info["so_luong"],
                    "vat_lieu": info["vat_lieu"],
                    "nguoi_lap": info["nguoi_lap"],
                    "ngay_ve": info["ngay_ve"],
                    "dung_sai_do_day": info["dung_sai_day"],
                    "dung_sai_kich_thuoc": info["dung_sai_khac"],
                    "kich_thuoc_tong_the": info["kich_thuoc"],
                    "trang_so": page_num + 1,
                    "doc_status": "pending_review",
                    "doc_id": doc_id,
                    "family_id": doc_info.get("family_id"),
                    "base_code": doc_info.get("base_code", ""),
                    "version_no": doc_info.get("version_no", 1),
                    "version_label": doc_info.get("version_label", ""),
                    "variant_code": doc_info.get("variant_code", "default"),
                    "variant_group": doc_info.get("variant_group", ""),
                    "lifecycle_status": doc_info.get("lifecycle_status", "draft"),
                    "review_status": doc_info.get("review_status", "pending_review"),
                    "is_current": doc_info.get("is_current", False),
                    "is_archived": doc_info.get("is_archived", False),
                    "supersedes_doc_id": doc_info.get("supersedes_doc_id"),
                }
                info['trang_so'] = page_num + 1
 
                # FIX #1: CHI insert metadata trang nay (khong xoa metadata cac trang khac)
                save_page_metadata(ten_file, thu_muc, info, doc_id=doc_id)
 
                all_chunks = []
                title_block = (
                    f"Thong tin tai lieu {ten_file}:\n"
                    f"- Ma chinh: {info.get('ma_chinh', [])}\n"
                    f"- Ma BTP: {info.get('ma_btp', [])}\n"
                    f"- Ma vat tu: {info.get('ma_vat_tu', [])}\n"
                    f"- Ma lien quan: {info.get('ma_lien_quan', [])}\n"
                    f"- Ma doi tuong tong hop: {info['ma_doi_tuong']}\n"
                    f"- Loai tai lieu: {info['loai_tai_lieu']}\n"
                    f"- Ten tai lieu/san pham: {info['ten_tai_lieu']}\n"
                    f"- Cong doan: {info['cong_doan']}\n"
                    f"- Vat lieu: {info['vat_lieu']}\n"
                    f"- So luong: {info['so_luong']}\n"
                    f"- Nguoi lap: {info['nguoi_lap']}\n"
                    f"- Ngay phat hanh: {info['ngay_ve']}\n"
                    f"- Dung sai do day vat lieu: {info['dung_sai_day']}\n"
                    f"- Dung sai cac kich thuoc khac: {info['dung_sai_khac']}\n"
                )
                if info['kich_thuoc']:
                    title_block += f"- Kich thuoc tong the: {info['kich_thuoc']}\n"
                all_chunks.append(Document(page_content=title_block, metadata={**metadata, "loai_du_lieu": "title_block"}))
 
                # FIX hieu nang: dung pdf_table_reader DA mo san thay vi mo lai moi trang
                markdown_tables = ""
                if pdf_table_reader is not None:
                    try:
                        page_plumber = pdf_table_reader.pages[page_num]
                        tables = page_plumber.extract_tables()
                        for table_idx, table in enumerate(tables):
                            # Parse BOM records and save to SQL
                            bom_records = extract_bom_records(table, table_idx=table_idx)
                            if bom_records:
                                save_bom_records(doc_id, page_num + 1, bom_records)
                                
                            for row in table:
                                cleaned_row = [str(cell).replace("\n", " ").strip() if cell else "" for cell in row]
                                markdown_tables += "| " + " | ".join(cleaned_row) + " |\n"
                            markdown_tables += "\n"
                    except Exception as e:
                        logger.warning(f"Khong the boc bang bieu {img_name}: {e}")
 
                if markdown_tables.strip():
                    table_content = f"Bang bieu tai lieu {ten_file} (Ma: {info['ma_doi_tuong']}):\n{markdown_tables}"
                    table_chunks = token_splitter.split_text(table_content)
                    for i, c in enumerate(table_chunks):
                        all_chunks.append(Document(page_content=c, metadata={**metadata, "loai_du_lieu": "bang_ke_vat_tu", "chunk_index": i}))
 
                if info['yckt']:
                    yckt_content = f"Yeu cau ky thuat tai lieu {ten_file} (Ma: {info['ma_doi_tuong']}):\n{info['yckt']}"
                    yckt_chunks = token_splitter.split_text(yckt_content)
                    for i, c in enumerate(yckt_chunks):
                        all_chunks.append(Document(page_content=c, metadata={**metadata, "loai_du_lieu": "yckt", "chunk_index": i}))
 
                if info['hdcv']:
                    hdcv_content = f"Huong dan cong viec tai lieu {ten_file} (Ma: {info['ma_doi_tuong']}):\n{info['hdcv']}"
                    hdcv_chunks = token_splitter.split_text(hdcv_content)
                    for i, c in enumerate(hdcv_chunks):
                        all_chunks.append(Document(page_content=c, metadata={**metadata, "loai_du_lieu": "hdcv", "chunk_index": i}))
 
                # Luoi an toan: luu raw text de giu moi thong tin chi tiet
                # (kich thuoc, goc, ban kinh, ghi chu) ma regex khong cover duoc
                if text.strip():
                    raw_content = f"Noi dung chi tiet trang {page_num+1} tai lieu {ten_file} (Ma: {info['ma_doi_tuong']}):\n{text.strip()}"
                    raw_chunks = token_splitter.split_text(raw_content)
                    for i, c in enumerate(raw_chunks):
                        all_chunks.append(Document(page_content=c, metadata={**metadata, "loai_du_lieu": "text", "chunk_index": i}))
 
                if image_summary.strip():
                    img_summary_content = f"Phan tich hinh anh tai lieu {ten_file} (Ma: {info['ma_doi_tuong']}):\n{image_summary}"
                    all_chunks.append(Document(page_content=img_summary_content, metadata={**metadata, "loai_du_lieu": "image_summary", **vision_metadata}))
 
                # FIX #3: GIU text goc cho LLM (noi_dung_goc) TRUOC khi tokenize ban dung cho BM25
                for chunk in all_chunks:
                    chunk.metadata["noi_dung_goc"] = chunk.page_content
                    chunk.page_content = tokenize_cached(chunk.page_content)
 
                # Document Versioning: Xoa vector cu cua file nay truoc khi add (chi xoa 1 lan o trang 1)
                if page_num == 0:
                    try:
                        _delete_vectors_for_file(ten_file, thu_muc)
                    except Exception as e:
                        logger.warning(f"Khong xoa duoc vector cu (bo qua, tiep tuc): {ten_file}: {e}")
 
                _add_docs_with_retry(all_chunks)
                report["total_chunks"] += len(all_chunks)
 
            except Exception as e:
                logger.error(f"Loi khi xu ly trang {page_num+1} cua {ten_file}: {e}", exc_info=True)
                report["failed_pages"].append(page_num+1)
                # FIX #2: KHONG dong doc o day nua (truoc day doc.close() trong except
                # khien cac trang sau cua cung file deu loi day chuyen).
 
    except Exception as e:
        logger.error(f"Loi doc file PDF {ten_file}: {e}", exc_info=True)
        report["status"] = "error"
        report["message"] = str(e)
    finally:
        # FIX #2: dong tai nguyen DUNG MOT LAN, ke ca khi thanh cong
        # (truoc day thieu doc.close() o nhanh thanh cong -> ro ri file handle).
        if pdf_table_reader is not None:
            try:
                pdf_table_reader.close()
            except Exception:
                pass
        if doc is not None:
            doc.close()
 
    if report["status"] == "success" and (report["failed_pages"] or report["total_chunks"] == 0):
        report["status"] = "error"

    if report["status"] == "error" and ROLLBACK_ON_INGEST_ERROR:
        try:
            _delete_vectors_for_file(ten_file, thu_muc)
            reset_document_metadata(ten_file, thu_muc)
            report["total_chunks"] = 0
            report["warnings"].append("Da rollback vector/metadata cua file nay vi ingest khong dat quality gate.")
        except Exception as e:
            report["warnings"].append(f"Rollback vector/metadata that bai: {e}")
            logger.warning(f"Rollback vector/metadata that bai cho {ten_file}: {e}")

    report["time_taken"] = round(time.time() - start_time, 2)
    message_parts = []
    if report["message"]:
        message_parts.append(report["message"])
    if report["failed_pages"]:
        message_parts.append(f"Cac trang loi/bo qua: {report['failed_pages']}")
    if report["vision_failed_pages"]:
        message_parts.append(f"Trang loi Gemini Vision/OCR: {report['vision_failed_pages']}")
    if report["metadata_llm_failed_pages"]:
        message_parts.append(f"Trang loi Gemini metadata fallback: {report['metadata_llm_failed_pages']}")
    if report["warnings"]:
        message_parts.append("Canh bao chat luong: " + " | ".join(report["warnings"][:5]))
    if not message_parts:
        message_parts.append(f"Da nap {report['total_chunks']} chunks tu {report['total_pages']} trang.")
    report["message"] = " ".join(message_parts)
    return report
 
def process_and_ingest_file(file_path, ten_file, thu_muc, vision_model=None, progress_callback=None):
    ext = os.path.splitext(file_path)[1].lower()
    start_time = time.time()
    report = {
        "status": "success",
        "ten_file": ten_file,
        "total_pages": 1,
        "total_chunks": 0,
        "failed_pages": [],
        "vision_failed_pages": [],
        "metadata_llm_failed_pages": [],
        "warnings": [],
        "time_taken": 0,
        "message": ""
    }
    try:
        if progress_callback:
            progress_callback(f"Dang doc noi dung file {ext}...")
        text_content, data_type = extract_text_from_supported_file(file_path, ten_file, vision_model)
        text_content = text_content.strip()
        if not text_content:
            raise ValueError("Khong trich xuat duoc noi dung co the tim kiem tu file nay.")
 
        warning_count_before_metadata = len(report["warnings"])
        info = extract_metadata_smart(
            text_content[:5000],
            ten_file,
            thu_muc,
            vision_model,
            quality_warnings=report["warnings"],
        )
        if len(report["warnings"]) > warning_count_before_metadata:
            report["metadata_llm_failed_pages"].append(1)
        if info.get("ten_tai_lieu") == "Khong ro":
            info["ten_tai_lieu"] = os.path.splitext(ten_file)[0]
 
        # Neu file tu hoc la hinh anh, copy no sang Data_Anh_Da_Tach de sau nay lam Ban ve can cu
        if ext in IMAGE_EXTENSIONS:
            import shutil
            safe_thu_muc = re.sub(r'[\\/*?:"<>|]', "", thu_muc) if thu_muc else ""
            base_name = os.path.splitext(ten_file)[0]
            if safe_thu_muc:
                img_name = f"{safe_thu_muc}_{base_name}_page1.png"
            else:
                img_name = f"{base_name}_page1.png"
            img_path = os.path.join(IMAGE_DIR, img_name)
            # Convert sang PNG va luu
            try:
                img = Image.open(file_path)
                img.save(img_path, format="PNG")
            except Exception as e:
                logger.warning(f"Loi khi copy anh tu hoc vao Data_Anh_Da_Tach: {e}")
 
        if ext != ".pdf" and info.get("loai_tai_lieu") == "Ban ve gia cong":
            type_map = {
                "bang_du_lieu": "Bang du lieu",
                "van_ban_word": "Tai lieu Word",
                "slide": "Tai lieu trinh chieu",
                "image_summary": "Tai lieu anh/OCR",
                "van_ban_html": "Tai lieu HTML",
                "du_lieu_json": "Du lieu JSON",
                "du_lieu_xml": "Du lieu XML",
                "van_ban": "Tai lieu van ban",
            }
            info["loai_tai_lieu"] = type_map.get(data_type, "Tai lieu tong hop")
 
        metadata = {
            "file_goc": ten_file,
            "phong_ban_quyen": thu_muc,
            "ma_doi_tuong": info["ma_doi_tuong"],
            "ma_chinh": info.get("ma_chinh", []),
            "ma_btp": info.get("ma_btp", []),
            "ma_vat_tu": info.get("ma_vat_tu", []),
            "ma_lien_quan": info.get("ma_lien_quan", []),
            "loai_tai_lieu": info["loai_tai_lieu"],
            "ten_san_pham": info["ten_tai_lieu"],
            "cong_doan": info["cong_doan"],
            "so_luong": info["so_luong"],
            "vat_lieu": info["vat_lieu"],
            "nguoi_lap": info["nguoi_lap"],
            "ngay_ve": info["ngay_ve"],
            "dung_sai_do_day": info["dung_sai_day"],
            "dung_sai_kich_thuoc": info["dung_sai_khac"],
            "kich_thuoc_tong_the": info["kich_thuoc"],
            "trang_so": 1,
            "dinh_dang_file": ext,
            "doc_status": "pending_review",
        }
        info["trang_so"] = 1
 
        # File 1 trang: dung wrapper save_document_metadata (reset + insert 1 lan)
        save_document_metadata(ten_file, thu_muc, info)
 
        all_chunks = []
        title_block = (
            f"Thong tin tai lieu {ten_file}:\n"
            f"- Ma chinh: {info.get('ma_chinh', [])}\n"
            f"- Ma BTP: {info.get('ma_btp', [])}\n"
            f"- Ma vat tu: {info.get('ma_vat_tu', [])}\n"
            f"- Ma lien quan: {info.get('ma_lien_quan', [])}\n"
            f"- Ma doi tuong tong hop: {info['ma_doi_tuong']}\n"
            f"- Loai tai lieu: {info['loai_tai_lieu']}\n"
            f"- Ten tai lieu/san pham: {info['ten_tai_lieu']}\n"
            f"- Cong doan/thu muc: {info['cong_doan']}\n"
            f"- Vat lieu: {info['vat_lieu']}\n"
            f"- Dinh dang file: {ext}\n"
        )
        all_chunks.append(Document(page_content=title_block, metadata={**metadata, "loai_du_lieu": "title_block"}))
 
        # Dung chung token_splitter da dinh nghia theo gioi han embedding.
        chunks = token_splitter.split_text(text_content)
        for i, chunk in enumerate(chunks):
            if chunk.strip():
                all_chunks.append(Document(
                    page_content=chunk,
                    metadata={**metadata, "loai_du_lieu": data_type, "chunk_index": i + 1}
                ))
 
        # FIX #3: GIU text goc cho LLM truoc khi tokenize cho BM25 (ap dung TOAN BO chunks)
        for chunk in all_chunks:
            chunk.metadata["noi_dung_goc"] = chunk.page_content
            chunk.page_content = tokenize_cached(chunk.page_content)
 
        if all_chunks:
            # Document Versioning: Xoa vector cu
            try:
                _delete_vectors_for_file(ten_file, thu_muc)
            except Exception as e:
                logger.warning(f"Khong xoa duoc vector cu (bo qua, tiep tuc): {ten_file}: {e}")
 
            _add_docs_with_retry(all_chunks)
            report["total_chunks"] += len(all_chunks)
 
    except Exception as e:
        logger.error(f"Loi doc file {ten_file}: {e}", exc_info=True)
        report["status"] = "error"
        report["message"] = str(e)
 
    if report["status"] == "success" and report["total_chunks"] == 0:
        report["status"] = "error"

    if report["status"] == "error" and ROLLBACK_ON_INGEST_ERROR:
        try:
            _delete_vectors_for_file(ten_file, thu_muc)
            reset_document_metadata(ten_file, thu_muc)
            report["total_chunks"] = 0
            report["warnings"].append("Da rollback vector/metadata cua file nay vi ingest khong dat quality gate.")
        except Exception as e:
            report["warnings"].append(f"Rollback vector/metadata that bai: {e}")
            logger.warning(f"Rollback vector/metadata that bai cho {ten_file}: {e}")

    report["time_taken"] = round(time.time() - start_time, 2)
    message_parts = []
    if report["message"]:
        message_parts.append(report["message"])
    if report["metadata_llm_failed_pages"]:
        message_parts.append(f"Trang loi Gemini metadata fallback: {report['metadata_llm_failed_pages']}")
    if report["warnings"]:
        message_parts.append("Canh bao chat luong: " + " | ".join(report["warnings"][:5]))
    if not message_parts:
        message_parts.append(f"Da nap {report['total_chunks']} chunks tu file {ext}.")
    report["message"] = " ".join(message_parts)
    return report
