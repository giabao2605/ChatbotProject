# -*- coding: utf-8 -*-
"""Auto-split tu ingestion/pdf_processor.py (P1.3). Giu nguyen logic goc; chi tach file + import."""

import os
import re
import json
import html
from PIL import Image

# cross-module (owned) imports
from mech_chatbot.ingestion.pdf.config import HTML_EXTENSIONS, IMAGE_EXTENSIONS, PRESENTATION_EXTENSIONS, SUPPORTED_LEARNING_EXTENSIONS, TABLE_EXTENSIONS, TEXT_EXTENSIONS, WORD_EXTENSIONS
from mech_chatbot.ingestion.pdf.vision import call_vision_model, format_vision_data, parse_vision_json


try:
    import pandas as pd
except ImportError:
    pd = None


try:
    import docx
except ImportError:
    docx = None


try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None


try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def _require_package(package, feature_name):
    if package is None:
        raise ImportError(f"Thieu thu vien de doc {feature_name}. Hay cai/cap nhat requirements.txt roi chay lai.")


def _read_text_file(file_path):
    encodings = ("utf-8-sig", "utf-8", "cp1258", "cp1252", "latin-1")
    last_error = None
    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError as exc:
            last_error = exc
    raise UnicodeDecodeError(
        last_error.encoding,
        last_error.object,
        last_error.start,
        last_error.end,
        "Khong doc duoc file bang cac encoding pho bien.",
    )


def _read_json_file(file_path):
    raw_text = _read_text_file(file_path)
    try:
        data = json.loads(raw_text)
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception:
        return raw_text


def _read_xml_file(file_path):
    raw_text = _read_text_file(file_path)
    try:
        import xml.etree.ElementTree as ET
        root = ET.fromstring(raw_text)
        extracted = "\n".join(t.strip() for t in root.itertext() if t and t.strip())
        return extracted or raw_text
    except Exception:
        return raw_text


def _read_html_file(file_path):
    raw_text = _read_text_file(file_path)
    if BeautifulSoup:
        soup = BeautifulSoup(raw_text, "html.parser")
        return soup.get_text("\n")
    without_tags = re.sub(r"<[^>]+>", " ", raw_text)
    return html.unescape(without_tags)


def _dataframe_to_text(df, title=None):
    if df is None or df.empty:
        return f"{title}\n(Bang rong)" if title else "(Bang rong)"
    df = df.fillna("")
    text = df.to_string(index=False)
    return f"{title}\n{text}" if title else text


def _read_table_file(file_path, ext):
    _require_package(pd, "bang tinh/CSV")
    if ext in {".csv", ".tsv"}:
        sep = "\t" if ext == ".tsv" else ","
        last_error = None
        for encoding in ("utf-8-sig", "utf-8", "cp1258", "cp1252", "latin-1"):
            try:
                df = pd.read_csv(file_path, sep=sep, encoding=encoding)
                return _dataframe_to_text(df, title=f"Bang du lieu tu {os.path.basename(file_path)}")
            except UnicodeDecodeError as exc:
                last_error = exc
        raise last_error
    sheets = pd.read_excel(file_path, sheet_name=None)
    sheet_texts = []
    for sheet_name, df in sheets.items():
        sheet_texts.append(_dataframe_to_text(df, title=f"Sheet: {sheet_name}"))
    return "\n\n---\n\n".join(sheet_texts)


def _read_word_file(file_path):
    _require_package(docx, "Word DOCX")
    document = docx.Document(file_path)
    parts = [para.text.strip() for para in document.paragraphs if para.text and para.text.strip()]
    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"Bang {table_index}:\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _read_presentation_file(file_path):
    _require_package(Presentation, "PowerPoint PPTX")
    prs = Presentation(file_path)
    slides = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    parts.append("Bang:\n" + "\n".join(rows))
        if parts:
            slides.append(f"Slide {slide_index}:\n" + "\n\n".join(parts))
    return "\n\n---\n\n".join(slides)


def _read_image_file(file_path, ten_file, vision_model):
    if not vision_model:
        raise ValueError("File ảnh cần PROXYLLM_API_KEY hợp lệ để GPT-5.4 Vision đọc nội dung/OCR.")
    image = Image.open(file_path)
    prompt = (
        f"Day la file anh '{ten_file}' duoc nap lam du lieu cho chatbot ky thuat. "
        "Hay OCR va tra ve ket qua DUOI DANG JSON voi schema sau:\n"
        "{\n"
        '  "document_codes": [],\n'
        '  "part_names": [],\n'
        '  "materials": [],\n'
        '  "dimensions": [],\n'
        '  "tolerances": [],\n'
        '  "technical_notes": [],\n'
        '  "bom_rows": [],\n'
        '  "uncertain_fields": []\n'
        "}\n"
        "Luon tra ve dung dinh dang JSON (khong kem text mo dau/ket thuc ngoai block ```json). "
        "Dien vao cac mang cac thong tin ky thuat tuong ung ban nhin thay trong anh."
    )
    response = call_vision_model(vision_model, prompt, image)
    vision_data = parse_vision_json(response.text)

    if vision_data:
        return format_vision_data(vision_data)

    return response.text or ""


def extract_text_from_supported_file(file_path, ten_file, vision_model=None):
    ext = os.path.splitext(file_path)[1].lower()
    if ext in {".json"}:
        return _read_json_file(file_path), "du_lieu_json"
    if ext in {".xml"}:
        return _read_xml_file(file_path), "du_lieu_xml"
    if ext in HTML_EXTENSIONS:
        return _read_html_file(file_path), "van_ban_html"
    if ext in TABLE_EXTENSIONS:
        return _read_table_file(file_path, ext), "bang_du_lieu"
    if ext in WORD_EXTENSIONS:
        return _read_word_file(file_path), "van_ban_word"
    if ext in PRESENTATION_EXTENSIONS:
        return _read_presentation_file(file_path), "slide"
    if ext in IMAGE_EXTENSIONS:
        return _read_image_file(file_path, ten_file, vision_model), "image_summary"
    if ext in TEXT_EXTENSIONS:
        return _read_text_file(file_path), "van_ban"
    supported = ", ".join(sorted(SUPPORTED_LEARNING_EXTENSIONS))
    raise ValueError(f"Dinh dang {ext or '(khong co duoi file)'} chua duoc ho tro. Cac dinh dang dang ho tro: {supported}")

__all__ = [
    '_require_package',
    'pd',
    'docx',
    'BeautifulSoup',
    'Presentation',
    '_read_text_file',
    '_read_json_file',
    '_read_xml_file',
    '_read_html_file',
    '_dataframe_to_text',
    '_read_table_file',
    '_read_word_file',
    '_read_presentation_file',
    '_read_image_file',
    'extract_text_from_supported_file',
]
